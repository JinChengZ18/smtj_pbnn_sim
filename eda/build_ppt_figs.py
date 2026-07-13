#!/usr/bin/env python3
"""Chapter-deck figure pipeline: slide position == chapter figure number.

Every numbered raster asset ``article/figs/Chapter0X_local_NN.png`` is
exported from slide ``NN`` of the hand-maintained deck
``article/ppt/Chapter0X_local.pptx``. The deck is the single place where
figure numbering and panel letters live; generator scripts only produce
raw, unnumbered plots (``figures/``, ``figures/panels/``, ``demo/figures/``).

What this script does per deck (idempotent):

  1. tag untagged hand-built slides with ``FIG:NN`` notes, in deck order
     following the manifest's ``hand_order``;
  2. rebuild the multi-panel AUTOFIG slides (letters added here, not in
     the plots) from ``figures/panels/`` exactly as before;
  3. insert missing figures: ``singles`` as one full-width raw image per
     slide, ``vector`` figures as position-holder slides carrying the
     current SVG-pipeline render (tagged ``FIG:NN:VECTOR``);
  4. refresh embedded images whose registered raw source changed
     (byte comparison);
  5. reorder slides so that slide index + 1 == figure number;
  6. render the deck via LibreOffice and export every non-vector slide,
     auto-cropped to content, to ``article/figs/<deck>_NN.png``.

Vector figures (circuit schematics / architecture SVGs) keep their
cairosvg-rendered assets; their deck slides only hold the position.

Run (Windows):  python eda/build_ppt_figs.py
"""
from __future__ import annotations

import re
import shutil
import subprocess
import sys
from pathlib import Path

import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from PIL import Image
import fitz

REPO = Path(__file__).resolve().parents[1]
PANELS = REPO / "figures" / "panels"
FIGS = REPO / "figures"
DEMO = REPO / "demo" / "figures"
PPT = REPO / "article" / "ppt"
OUT = REPO / "article" / "figs"
SOFFICE = r"C:\Program Files\LibreOffice\program\soffice.exe"
MARGIN, GAP, TOP = 0.25, 0.14, 0.55
LETTER_DX, LETTER_DY = 0.05, 0.02
EXPORT_DPI = 419
CROP_PAD = 8  # px kept around the auto-cropped content

MANIFESTS = {
    "Chapter01_local": {
        "n_figs": 3,
        "hand_order": [1, 2, 3],
        "autofigs": {},
        "singles": {},
        "vector": {},
        "refresh": {},
    },
    "Chapter04_local": {
        "n_figs": 25,
        "hand_order": [1, 2, 3, 4, 5, 6, 7, 8, 9, 10, 11, 12, 13],
        "autofigs": {15: ("ch04_15", "abc"), 16: ("ch04_16", "ab"),
                     17: ("ch04_17", "abc"), 22: ("ch04_22", "ab")},
        # composed multi-panel plots that ship as one image: letters overlaid
        # at the detected axes corners (fig 19 = 3 transient-waveform panels)
        "overlays": {19: (FIGS / "waveforms_3ops.png", "abc", (1, 3)),
                     23: (FIGS / "structure_consistency.png", "ab", (1, 2)),
                     24: (FIGS / "30_temperature_selfconsistency.png", "abc", (1, 3)),
                     25: (FIGS / "32_replay_column_cosim.png", "ab", (1, 2))},
        "singles": {20: FIGS / "mode_pipeline.png",
                    21: FIGS / "dual_model_consistency.png"},
        "vector": {14: OUT / "Chapter04_local_14.png",
                   18: OUT / "Chapter04_local_18.png"},
        "refresh": {13: FIGS / "13a_training_energy_breakdown.png",
                    20: FIGS / "mode_pipeline.png",
                    21: FIGS / "dual_model_consistency.png"},
    },
    "Chapter05_local": {
        "n_figs": 11,
        "hand_order": [1, 2, 3, 4, 5, 6, 7],
        "autofigs": {9: ("ch05_09", "abc")},
        "overlays": {11: (FIGS / "31_rc_counting_readout.png", "ab", (1, 2))},
        "singles": {},
        "vector": {8: OUT / "Chapter05_local_08.png",
                   10: OUT / "Chapter05_local_10.png"},
        "refresh": {6: FIGS / "18_rc_benchmarks.png"},
    },
}

# Appendix decks are fully generated: one slide per figure, the composed raw
# plot placed full-width with panel letters overlaid at the detected axes-frame
# corners (letters live in the deck, never baked into the raw plot). Each entry
# is fig_number -> (raw source, panel letters, subplot grid (rows, cols)).
APPENDIX = {
    "AppendixB": {
        1: (FIGS / "05a_fashion_mnist_training_curves.png", "ab", (1, 2)),
        2: (FIGS / "05a_cifar10_training_curves.png", "ab", (1, 2)),
        3: (FIGS / "06a_fashion_mnist_sweep_T.png", "ab", (1, 2)),
        4: (FIGS / "06a_cifar10_sweep_T.png", "ab", (1, 2)),
    },
    "AppendixC": {
        1: (FIGS / "21_seed_independence.png", "abcd", (2, 2)),
    },
    "AppendixD": {
        8: (FIGS / "sar_capdac_switching.png", "ab", (1, 2)),
    },
}

LABEL_RE = re.compile(r"^\([a-z]\)$")
# panel letter sits this far up-and-left of each subplot's axes-frame corner
LABEL_DX_IN, LABEL_DY_IN = 0.34, 0.42


# ---------------------------------------------------------------- helpers
def note_of(slide) -> str:
    return (slide.notes_slide.notes_text_frame.text.strip()
            if slide.has_notes_slide else "")


def set_note(slide, text: str) -> None:
    slide.notes_slide.notes_text_frame.text = text


def fig_of(slide):
    n = note_of(slide)
    if n.startswith("FIG:"):
        return int(n.split(":")[1])
    return None


def walk_pics(shapes):
    for sh in shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from walk_pics(sh.shapes)
        elif sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            yield sh


def walk_shapes(shapes):
    for sh in shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from walk_shapes(sh.shapes)
        else:
            yield sh


def normalize_panel_labels(slide) -> int:
    """Harden hand-placed panel-letter textboxes so the closing ``)`` can never
    wrap to a second line: merge any split runs into one and disable word-wrap
    (matching the script-built AUTOFIG labels). Existing frame insets are kept
    so labels that already render correctly do not shift. Font size/bold/name/
    color are preserved; unset bold defaults to True. Returns the number of
    labels normalized."""
    fixed = 0
    for sh in walk_shapes(slide.shapes):
        if not sh.has_text_frame:
            continue
        tf = sh.text_frame
        if not LABEL_RE.match(tf.text.strip()):
            continue
        text = tf.text.strip()
        size = bold = name = color = None
        for p in tf.paragraphs:
            for r in p.runs:
                f = r.font
                if size is None and f.size is not None:
                    size = f.size
                if bold is None and f.bold is not None:
                    bold = f.bold
                if name is None and f.name is not None:
                    name = f.name
                try:
                    if color is None and f.color and f.color.type is not None:
                        color = f.color.rgb
                except (AttributeError, TypeError):
                    pass
        tf.word_wrap = False
        p0 = tf.paragraphs[0]
        for extra in tf.paragraphs[1:]:
            extra._p.getparent().remove(extra._p)
        for r in list(p0.runs):
            r._r.getparent().remove(r._r)
        run = p0.add_run()
        run.text = text
        run.font.bold = True if bold is None else bold
        if size is not None:
            run.font.size = size
        if name is not None:
            run.font.name = name
        if color is not None:
            run.font.color.rgb = color
        fixed += 1
    return fixed


def frame_corners(src: Path, rows: int, cols: int):
    """Top-left corner (xfrac, yfrac) of each subplot's axes frame, row-major.

    The composed plot is split into a rows x cols grid of cells; within each
    cell the axes frame is the first column/row carrying a long dark run (the
    spine). Used to anchor panel letters over already-composed figures."""
    im = Image.open(src)
    W, H = im.size
    dark = np.asarray(im.convert("L")) < 120
    cw, ch = W // cols, H // rows
    out = []
    for r in range(rows):
        for c in range(cols):
            x0, y0 = c * cw, r * ch
            sub = dark[y0:y0 + ch, x0:x0 + cw]
            hh, ww = sub.shape
            col_run, row_run = sub.sum(axis=0), sub.sum(axis=1)
            lc = next((i for i in range(ww) if col_run[i] > 0.45 * hh), 0)
            tr = next((i for i in range(hh) if row_run[i] > 0.45 * ww), 0)
            out.append(((x0 + lc) / W, (y0 + tr) / H))
    return out


def add_white_bg(slide, prs) -> None:
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    bg.line.fill.background(); bg.shadow.inherit = False
    # Drop the theme shape style: its effectRef applies a drop shadow that
    # survives shadow.inherit=False and shades the exported bottom/right
    # corners, defeating the crop-to-content step.
    style = bg._element.find(qn("p:style"))
    if style is not None:
        bg._element.remove(style)


def blank_slide(prs):
    layout = min(prs.slide_layouts, key=lambda l: len(l.placeholders))
    slide = prs.slides.add_slide(layout)
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)
    add_white_bg(slide, prs)
    return slide


def clear_shapes(slide) -> None:
    for sh in list(slide.shapes):
        sh._element.getparent().remove(sh._element)


def _letter_box(slide, lx: float, ly: float, ch: str) -> None:
    tb = slide.shapes.add_textbox(Inches(lx), Inches(ly),
                                  Inches(0.6), Inches(0.30))
    tf = tb.text_frame; tf.word_wrap = False
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    r = tf.paragraphs[0].add_run()
    r.text = f"({ch})"
    r.font.bold = True; r.font.size = Pt(14)
    r.font.color.rgb = RGBColor(0x20, 0x20, 0x20)


def add_single(prs, img: Path, tag: str):
    slide = blank_slide(prs)
    im = Image.open(img)
    pw = 10.0 - 2 * MARGIN
    ph = pw * im.height / im.width
    if ph > 6.8:                       # keep tall figures inside the page
        ph = 6.8; pw = ph * im.width / im.height
    x = (10.0 - pw) / 2
    slide.shapes.add_picture(str(img), Inches(x), Inches(TOP), width=Inches(pw))
    set_note(slide, tag)
    return slide


def populate_autofig(slide, stem: str, letters: str) -> None:
    """Fill an (already blank) slide with separate panel PNGs side by side,
    each carrying its panel letter (letters added here, never in the plots)."""
    paths = [PANELS / f"{stem}_{c}.png" for c in letters]
    n = len(paths)
    avail = 10.0 - 2 * MARGIN - (n - 1) * GAP
    pw = avail / n
    for i, path in enumerate(paths):
        x = MARGIN + i * (pw + GAP)
        slide.shapes.add_picture(str(path), Inches(x), Inches(TOP),
                                 width=Inches(pw))
        if n > 1:
            _letter_box(slide, x + LETTER_DX, TOP + LETTER_DY, letters[i])


def populate_overlay(slide, src: Path, letters: str, grid) -> None:
    """Place one composed multi-panel plot full-width on ``slide`` and overlay
    each panel letter at that subplot's detected axes-frame corner. Used for
    figures that ship as a single image (no separate panel PNGs)."""
    rows, cols = grid
    im = Image.open(src)
    pw = 10.0 - 2 * MARGIN
    ph = pw * im.height / im.width
    if ph > 6.8:
        ph = 6.8; pw = ph * im.width / im.height
    x = (10.0 - pw) / 2
    slide.shapes.add_picture(str(src), Inches(x), Inches(TOP), width=Inches(pw))
    for (xf, yf), ch in zip(frame_corners(src, rows, cols), letters):
        _letter_box(slide, max(0.02, x + xf * pw - LABEL_DX_IN),
                    max(0.02, TOP + yf * ph - LABEL_DY_IN), ch)


def add_autofig(prs, stem: str, letters: str, tag: str):
    slide = blank_slide(prs)
    populate_autofig(slide, stem, letters)
    set_note(slide, tag)
    return slide


def refresh_pic(slide, src: Path) -> bool:
    """Swap the slide's single largest picture for ``src`` if bytes differ."""
    pics = list(walk_pics(slide.shapes))
    if not pics:
        return False
    pic = max(pics, key=lambda p: p.width * p.height)
    part = pic.part.related_part(pic._element.blip_rId)
    new = src.read_bytes()
    if part.blob == new:
        return False
    part._blob = new
    im = Image.open(src)
    pic.height = int(pic.width * im.height / im.width)  # keep true aspect
    return True


def crop_content(im: Image.Image) -> Image.Image:
    """Crop a rendered page to its non-white content, keeping CROP_PAD px."""
    bbox = Image.eval(im.convert("L"),
                      lambda p: 0 if p > 250 else 255).getbbox()
    if bbox:
        x0, y0, x1, y1 = bbox
        im = im.crop((max(0, x0 - CROP_PAD), max(0, y0 - CROP_PAD),
                      min(im.width, x1 + CROP_PAD), min(im.height, y1 + CROP_PAD)))
    return im


def reorder(prs, order_rids):
    lst = prs.slides._sldIdLst
    id_by_r = {sid.get(qn("r:id")): sid for sid in list(lst)}
    for sid in list(lst):
        lst.remove(sid)
    for rid in order_rids:
        lst.append(id_by_r[rid])


# ---------------------------------------------------------------- main
def process(name: str, man: dict) -> None:
    deck = PPT / f"{name}.pptx"
    bak = deck.with_suffix(".pptx.bak")
    if not bak.exists():
        shutil.copy2(deck, bak)
    prs = Presentation(str(deck))

    # 1. (re)build every AUTOFIG / OVERLAY slot in place: reuse the existing
    #    FIG:NN slide (or a legacy AUTOFIG:stem placeholder) by clearing and
    #    repopulating it so regenerated panels always propagate; only add a
    #    fresh slide when the slot is absent. Never drop+add — that would leave
    #    an orphaned slide part and corrupt the package (duplicate partname).
    def slot_slide(nn, legacy=None):
        for s in prs.slides:
            if fig_of(s) == nn or (legacy and note_of(s) == legacy):
                return s
        return None

    def rebuild_slot(nn, fill, legacy=None):
        slide = slot_slide(nn, legacy)
        if slide is None:
            slide = blank_slide(prs)
        else:
            clear_shapes(slide); add_white_bg(slide, prs)
        fill(slide)
        set_note(slide, f"FIG:{nn:02d}")

    for nn, (stem, letters) in man["autofigs"].items():
        rebuild_slot(nn, lambda s, st=stem, le=letters: populate_autofig(s, st, le),
                     legacy=f"AUTOFIG:{stem}")
    for nn, (src, letters, grid) in man.get("overlays", {}).items():
        rebuild_slot(nn, lambda s, sr=src, le=letters, g=grid:
                     populate_overlay(s, sr, le, g))

    # 2. tag untagged hand slides in deck order per manifest
    hand_iter = iter(man["hand_order"])
    for slide in prs.slides:
        if fig_of(slide) is None and not note_of(slide).startswith("FIG:"):
            try:
                set_note(slide, f"FIG:{next(hand_iter):02d}")
            except StopIteration:
                raise RuntimeError(f"{name}: more untagged slides than "
                                   f"hand_order entries")

    have = {fig_of(s) for s in prs.slides if fig_of(s) is not None}

    # 3. insert singles / vector holders (skip slots already present so
    #    re-runs stay idempotent)
    for nn, img in man["singles"].items():
        if nn not in have:
            add_single(prs, img, f"FIG:{nn:02d}")
    for nn, img in man["vector"].items():
        if nn not in have:
            add_single(prs, img, f"FIG:{nn:02d}:VECTOR")

    # 4. refresh registered embeds
    for nn, src in man["refresh"].items():
        for slide in prs.slides:
            if fig_of(slide) == nn and refresh_pic(slide, src):
                print(f"  {name} fig {nn:02d}: embedded image refreshed")

    # 5. order slides by figure number
    pairs = []
    for sid, slide in zip(list(prs.slides._sldIdLst), list(prs.slides)):
        nn = fig_of(slide)
        if nn is None:
            raise RuntimeError(f"{name}: slide without FIG tag")
        pairs.append((nn, sid.get(qn("r:id"))))
    pairs.sort()
    nums = [nn for nn, _ in pairs]
    if nums != list(range(1, man["n_figs"] + 1)):
        raise RuntimeError(f"{name}: figure set {nums} != 1..{man['n_figs']}")
    reorder(prs, [rid for _, rid in pairs])

    # 5b. harden every hand-placed panel letter so ')' can never wrap
    nfix = sum(normalize_panel_labels(s) for s in prs.slides)
    if nfix:
        print(f"  {name}: normalized {nfix} panel-letter label(s)")
    prs.save(str(deck))

    # 6. render + export
    subprocess.run([SOFFICE, "--headless", "--convert-to", "pdf",
                    "--outdir", str(PPT), str(deck)],
                   check=True, capture_output=True)
    pdf = fitz.open(str(PPT / f"{name}.pdf"))
    vector_set = set(man["vector"])
    for nn in range(1, man["n_figs"] + 1):
        if nn in vector_set:
            continue                      # asset stays with the SVG pipeline
        page = pdf.load_page(nn - 1)
        pix = page.get_pixmap(dpi=EXPORT_DPI)
        im = crop_content(Image.frombytes("RGB", (pix.width, pix.height),
                                          pix.samples))
        out = OUT / f"{name}_{nn:02d}.png"
        im.save(out)
        print(f"  wrote {out.relative_to(REPO)}  ({im.width}x{im.height})")
    pdf.close()


def build_appendix(name: str, overlays: dict) -> None:
    """Generate an appendix deck from scratch: one slide per figure, the raw
    composed plot full-width with panel letters overlaid at the detected
    axes-frame corners, then export ``<name>_NN.png`` for every figure."""
    deck = PPT / f"{name}.pptx"
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(10), Inches(7.5)
    figs = sorted(overlays)
    for nn in figs:
        src, letters, (rows, cols) = overlays[nn]
        slide = blank_slide(prs)
        im = Image.open(src)
        pw = 10.0 - 2 * MARGIN
        ph = pw * im.height / im.width
        if ph > 6.8:
            ph = 6.8; pw = ph * im.width / im.height
        x = (10.0 - pw) / 2
        slide.shapes.add_picture(str(src), Inches(x), Inches(TOP),
                                 width=Inches(pw))
        for (xf, yf), ch in zip(frame_corners(src, rows, cols), letters):
            lx = max(0.02, x + xf * pw - LABEL_DX_IN)
            ly = max(0.02, TOP + yf * ph - LABEL_DY_IN)
            tb = slide.shapes.add_textbox(Inches(lx), Inches(ly),
                                          Inches(0.6), Inches(0.30))
            tf = tb.text_frame; tf.word_wrap = False
            tf.margin_left = tf.margin_right = 0
            tf.margin_top = tf.margin_bottom = 0
            r = tf.paragraphs[0].add_run()
            r.text = f"({ch})"
            r.font.bold = True; r.font.size = Pt(14)
            r.font.color.rgb = RGBColor(0x20, 0x20, 0x20)
        set_note(slide, f"FIG:{nn:02d}")
    prs.save(str(deck))

    subprocess.run([SOFFICE, "--headless", "--convert-to", "pdf",
                    "--outdir", str(PPT), str(deck)],
                   check=True, capture_output=True)
    pdf = fitz.open(str(PPT / f"{name}.pdf"))
    for i, nn in enumerate(figs):
        pix = pdf.load_page(i).get_pixmap(dpi=EXPORT_DPI)
        im = crop_content(Image.frombytes("RGB", (pix.width, pix.height),
                                          pix.samples))
        out = OUT / f"{name}_{nn:02d}.png"
        im.save(out)
        print(f"  wrote {out.relative_to(REPO)}  ({im.width}x{im.height})")
    pdf.close()


def main(argv=None) -> None:
    """Process every deck, or only those named on the command line.

    ``python eda/build_ppt_figs.py``  -> all chapter + appendix decks
    ``python eda/build_ppt_figs.py Chapter05_local AppendixB``  -> just those
    """
    sel = set(argv) if argv else None
    for name, man in MANIFESTS.items():
        if sel and name not in sel:
            continue
        print(f"== {name}")
        process(name, man)
    for name, overlays in APPENDIX.items():
        if sel and name not in sel:
            continue
        print(f"== {name} (appendix)")
        build_appendix(name, overlays)


if __name__ == "__main__":
    main(sys.argv[1:] or None)
